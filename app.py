# app.py
import uuid
import os
from io import BytesIO
import pandas as pd
import streamlit as st

# --- Your services (keep your existing modules) ---
from services.chemical_profile import extract_chemical_profile
from services.excel_loader import read_file_dynamic
from services.extract_chemical_profile import extract_api_from_orange_book
from services.herbalengine import compute_herbal_opportunity  # (your module name as provided)
from services.overview_service import build_executive_overview
from services.cost_tool import formulation_cost_estimate
from services.llm_setup import init_llms
from services.generic_entry import earliest_generic_entry
from services.recommendation import recommendation_flag
from services.competition import competitive_density_score
from services.formulation_intelligence import (
    compute_formulation_risk,
    infer_formulation_type,
    is_high_value_complex,
)

from services.orange_book import (
    build_ege_table_cached,
    lookup_by_drug,
    lookup_by_patent,
    get_drug_dropdown_list,
    get_patents_for_drug,
    load_orange_book_products,
    load_orange_book_patents,
    load_orange_book_exclusivities,
)

from services.market_insights import (
    find_all_generic_drugs,
    find_expiring_patents,
)

from services.opportunities import rank_top_orange_book_opportunities

# --- PPT Export ---
from pptx import Presentation
from pptx.util import Inches, Pt


# =====================================================
# Helpers: bullet formatting for PPT
# =====================================================
def _add_bullets_to_textframe(tf, text: str):
    """
    Auto-detect bullets from AI text and render nicely:
    - Lines beginning with -, *, •, or numbered "1." become bullets
    - Blank lines are skipped
    """
    lines = [ln.strip() for ln in (text or "").splitlines()]
    lines = [ln for ln in lines if ln]

    if not lines:
        tf.text = "—"
        return

    # first paragraph
    first = True
    for ln in lines:
        is_bullet = False

        if ln.startswith(("-", "*", "•")):
            ln = ln.lstrip("-*•").strip()
            is_bullet = True
        else:
            # numbered bullet like "1) ..." or "1. ..."
            if len(ln) >= 3 and ln[0].isdigit() and (ln[1] in [".", ")"]):
                # handle 10., 12) etc
                # simple parse: split at first space after punctuation
                # keep content after "1." / "1)"
                parts = ln.split(maxsplit=1)
                if len(parts) == 2 and (parts[0].endswith(".") or parts[0].endswith(")")):
                    ln = parts[1].strip()
                    is_bullet = True

        if first:
            p = tf.paragraphs[0]
            p.text = ln
            p.level = 0 if is_bullet else 0
            first = False
        else:
            p = tf.add_paragraph()
            p.text = ln
            p.level = 0 if is_bullet else 0


# =====================================================
# PPT Export: builds a deck from tab outputs
# =====================================================
def export_report_to_ppt(
    commercial_mode: str,
    selected_drug: str,
    overview_text: str,
    ege_text: str,
    strategy_text: str,
    ai_text: str,
) -> BytesIO:
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Pharma Commercial Decision Engine"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Commercial Intelligence Report\nMode: {commercial_mode}\nDrug: {selected_drug or '—'}"

    # utility: add a title+content slide
    def add_text_slide(title: str, body: str):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.clear()
        _add_bullets_to_textframe(tf, body)

    add_text_slide("Executive Overview", overview_text)
    add_text_slide("EGE (Generic Entry) Summary", ege_text)
    add_text_slide("Strategy Summary", strategy_text)
    add_text_slide("AI Analysis", ai_text)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# =====================================================
# Dynamic Herbal Dataset Loader (3,000+ Plants)
# =====================================================

@st.cache_data(show_spinner=True)
def load_herbal_master_list():
    try:
        df = pd.read_csv("data/global_medicinal_plants_unique.csv")

        # Prefer Scientific Name but include Common Name
        if "Scientific_Name" in df.columns:
            plants = df["Common_Name"].dropna().unique().tolist()
        else:
            plants = df.iloc[:, 0].dropna().unique().tolist()

        plants = sorted(plants)

        return plants

    except Exception as e:
        st.error(f"Failed to load medicinal plant dataset: {e}")
        return []

def reload_app_after_upload():
    st.cache_data.clear()
    st.cache_resource.clear()

    for key in [
        "selected_drug",
        "selected_patent",
        "selected_opportunity_drug",
        "ege_master_df",
        "top20_opportunities",
        "top20_computed",
        "ai_summary",
        "board_summary",
        "run_id",
        "next_plan_text",
        "next_plan_run_id",
        "overview_snapshot",
        "ege_snapshot",
        "strategy_snapshot",
        "ai_snapshot",
    ]:
        st.session_state.pop(key, None)

    st.rerun()


# =====================================================
# Optional Meta-Agent
# =====================================================
try:
    from Agent.meta_agent import build_agent
except Exception:
    build_agent = None


# =====================================================
# Page setup
# =====================================================
st.set_page_config(page_title="Pharma Commercial Decision Engine", layout="wide")
st.title("🧪 Pharma Commercial Decision Engine")
st.caption("Orange Book–driven generic drug opportunity intelligence • Not legal or financial advice")


# =====================================================
# Session defaults
# =====================================================
DEFAULTS = {
    "commercial_mode": "FDA Orange Book (Rx Generics)",
    "selection_mode": "FDA Orange Book All",
    "selected_drug": "",
    "selected_patent": "",
    "selected_opportunity_drug": "",
    "herbal_input": "",
    "herbal_dropdown": "",
    "run_id": None,
    "ai_summary": None,
    "board_summary": None,
    "route": None,
    "dosage_form": None,
    "cost": None,
    "months_horizon": 12,
    "meta_agent": None,
    "meta_agent_model": None,
    "ege_master_df": pd.DataFrame(),
    "top20_opportunities": None,
    "top20_computed": False,
    "next_plan_text": None,
    "next_plan_run_id": None,
    "next_plan_style": "Corporate BD / IC Memo",
    # snapshots for PPT
    "overview_snapshot": "",
    "ege_snapshot": "",
    "strategy_snapshot": "",
    "ai_snapshot": "",
    "ai_summary": None,
    "ai_chat_history": [],
}
for k, v in DEFAULTS.items():
    st.session_state.setdefault(k, v)


def on_drug_change():
    # reset analysis-only values when drug changes
    for k in ["selected_patent", "run_id", "ai_summary", "board_summary", "cost", "next_plan_text", "next_plan_run_id"]:
        st.session_state[k] = None if k != "selected_patent" else ""


# =====================================================
# Load LLMs
# =====================================================
llm_openai, llm_groq, llm_ollama, llm_labels = init_llms()
if not llm_labels:
    st.error("No LLMs available.")
    st.stop()

model_choice = st.sidebar.selectbox("AI Model", llm_labels, key="model_choice")
llm = llm_openai if model_choice.startswith("OpenAI") else llm_groq if model_choice.startswith("Groq") else llm_ollama

if build_agent and (st.session_state.meta_agent is None or st.session_state.meta_agent_model != model_choice):
    try:
        st.session_state.meta_agent = build_agent(llm)
        st.session_state.meta_agent_model = model_choice
    except Exception:
        st.session_state.meta_agent = None


# =====================================================
# Cached opportunity dropdown
# =====================================================
@st.cache_data(show_spinner=False)
def build_top_opportunity_dropdown(max_drugs=1500, top_n=1000):
    products = load_orange_book_products()
    patents = load_orange_book_patents()
    exclus = load_orange_book_exclusivities()

    ranked = rank_top_orange_book_opportunities(
        products=products,
        patents=patents,
        exclus=exclus,
        top_n=top_n,
        max_drugs=max_drugs,
    )
    if ranked is None or ranked.empty:
        return []
    col = "Ingredient" if "Ingredient" in ranked.columns else "Drug"
    return ranked[col].dropna().unique().tolist()


# =====================================================
# Sidebar — Commercial Domain
# =====================================================
st.sidebar.subheader("Commercial Domain")
commercial_mode = st.sidebar.radio(
    "Select Domain",
    ["FDA Orange Book (Rx Generics)", "Herbal / Botanical Drugs"],
    key="commercial_mode",
)
st.sidebar.divider()


# =====================================================
# Sidebar — Selection
# (IMPORTANT: only render ONE patent widget, with one key)
# =====================================================
if commercial_mode == "FDA Orange Book (Rx Generics)":
    st.sidebar.subheader("Selection Mode")
    selection_mode = st.sidebar.radio(
        "Choose selection source",
        ["FDA Orange Book All", "🔥 Top Opportunities"],
        key="selection_mode",
    )

    st.sidebar.divider()

    if selection_mode == "FDA Orange Book All":
        drug_list = get_drug_dropdown_list()
        st.sidebar.selectbox(
            "Drug (FDA Orange Book)",
            [""] + drug_list,
            key="selected_drug",
            on_change=on_drug_change,
        )
    else:
        with st.spinner("Loading top opportunities…"):
            opp_list = build_top_opportunity_dropdown()

        st.sidebar.selectbox(
            "🔥 Top-1000 Opportunity Drugs",
            [""] + opp_list,
            key="selected_opportunity_drug",
        )

        # Sync selection to selected_drug
        if st.session_state.selected_opportunity_drug and st.session_state.selected_drug != st.session_state.selected_opportunity_drug:
            st.session_state.selected_drug = st.session_state.selected_opportunity_drug
            on_drug_change()

    # Patent options for FDA mode only
    patent_options = get_patents_for_drug(st.session_state.selected_drug) if st.session_state.selected_drug else []
else:
    st.sidebar.subheader("Herbal Drug Selection")

    herbal_list = load_herbal_master_list()

    st.sidebar.text_input(
        "Type Herbal Drug Name (ANY)",
        key="herbal_input",
        placeholder="e.g. Berberine, Lion's Mane, Shilajit…",
    )
    st.sidebar.caption("OR choose from common list (optional)")
    st.sidebar.selectbox(
        "Common Herbal Drugs",
        [""] + herbal_list,
        key="herbal_dropdown",
    )

    # Apply typed input first, fallback to dropdown
    if st.session_state.herbal_input.strip():
        st.session_state.selected_drug = st.session_state.herbal_input.strip()
    elif st.session_state.herbal_dropdown:
        st.session_state.selected_drug = st.session_state.herbal_dropdown
    else:
        # keep what user had
        pass

    patent_options = []  # not applicable


# Render ONE patent widget total (no duplicate keys)
st.sidebar.selectbox(
    "Patent (optional)",
    [""] + patent_options,
    key="selected_patent",
    disabled=(commercial_mode != "FDA Orange Book (Rx Generics)"),
)

st.sidebar.divider()


# =====================================================
# Analyze Button
# =====================================================
if st.sidebar.button(
    "🔍 Analyze",
    type="primary",
    disabled=not bool(st.session_state.selected_drug),
    key="analyze_btn",
):
    st.session_state.cost = formulation_cost_estimate(
        st.session_state.selected_drug,
        st.session_state.route or "",
        st.session_state.dosage_form or "",
    )
    st.session_state.run_id = str(uuid.uuid4())
    st.session_state.ai_summary = None
    st.session_state.board_summary = None
    st.session_state.next_plan_text = None
    st.session_state.next_plan_run_id = None


# =====================================================
# Core data (SAFE: avoid heavy work when herbal mode)
# =====================================================
selected_drug = st.session_state.selected_drug
ob_data = None
products = patents = exclus = pd.DataFrame()
generic_df = expiring_df = pd.DataFrame()

if commercial_mode == "FDA Orange Book (Rx Generics)" and selected_drug:
    try:
        ob_data = lookup_by_patent(st.session_state.selected_patent) if st.session_state.selected_patent else lookup_by_drug(selected_drug)
    except Exception as e:
        st.error(f"Failed to lookup Orange Book data: {e}")
        ob_data = None

    # load portfolio data
    try:
        products = load_orange_book_products()
        patents = load_orange_book_patents()
        exclus = load_orange_book_exclusivities()
        generic_df = find_all_generic_drugs(products, patents, exclus)
        expiring_df = pd.DataFrame(find_expiring_patents(patents, st.session_state.months_horizon))
    except Exception as e:
        st.error(f"Failed to load Orange Book datasets: {e}")
        products = patents = exclus = pd.DataFrame()
        generic_df = expiring_df = pd.DataFrame()


def build_next_level_plan_text(
    llm,
    style: str,
    drug: str,
    formulation_type: str,
    formulation_risk_level: str,
    formulation_risk_score: int | None,
    competition_level: str,
    competition_count: int | None,
    ege_years: float | None,
    cmc_cost: dict | None,
) -> str:
    ege_text = "Open now" if ege_years is None else f"{ege_years:.1f} years"
    cost_text = "Unknown" if not cmc_cost else f"${cmc_cost.get('low')}–${cmc_cost.get('high')}M"

    style_map = {
        "VC Pitch": "VC pitch: concise, value-creation, milestone-driven, speed-focused.",
        "PE Diligence": "PE diligence: risk controls, cost certainty, execution gates, downside protection.",
        "Corporate BD / IC Memo": "Corporate BD / IC memo: strategic fit, regulatory path, operational readiness, partnership angles.",
    }
    style_instr = style_map.get(style, style_map["Corporate BD / IC Memo"])

    prompt = f"""
You are a US pharma commercialization strategy lead.

Write a section titled:
"Next-Level Execution Plan (Commercial-Scale Readiness)"

Audience style:
{style_instr}

Context (do NOT invent facts or numbers):
- Drug: {drug}
- Formulation type: {formulation_type}
- Formulation risk: {formulation_risk_level}{f" ({formulation_risk_score}/100)" if formulation_risk_score is not None else ""}
- Competition: {competition_level}{f" ({competition_count})" if competition_count is not None else ""}
- Earliest Generic Entry (EGE): {ege_text}
- Estimated CMC cost: {cost_text}

Required structure (use numbered headings 1–5):
1) Formulation stability & scale-up validation
2) Supply-chain lock-in & cost certainty
3) Market access & pricing optimization
4) Regulatory & exclusivity leverage (mention 180-day exclusivity only as feasibility evaluation, not a claim)
5) Rapid-launch operational plan

Constraints:
- No legal advice
- No fabricated dates, market sizes, or pricing numbers
- Keep it 120–200 words total
- Use crisp, executive language
Return markdown only (no code fences).
"""
    return llm.invoke(prompt).content.strip()


# =====================================================
# Tabs
# (NO st.stop() anywhere in tabs — use if/else only)
# =====================================================
overview_tab, ege_tab, strategy_tab, ai_tab, data_upload = st.tabs(
    ["🏠 Executive Overview", "⏳ Generic Entry (EGE)", "📊 Strategy", "🧠 AI Analysis", "📂 Data Upload"]
)

# =====================================================
# OVERVIEW TAB
# =====================================================
with overview_tab:
    st.subheader("Executive Overview")

    overview_snapshot_lines = []

    if commercial_mode == "Herbal / Botanical Drugs":
        if not selected_drug:
            st.info("Type a herbal name or pick from the list.")
            overview_snapshot_lines.append("No herbal selected.")
        else:
            herbal_info = compute_herbal_opportunity(selected_drug)

            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Herbal Drug", selected_drug)
            c2.metric("Demand Score", herbal_info.get("Demand Score", "—"))
            c3.metric("Formulation Complexity", herbal_info.get("Formulation Complexity", "—"))
            c4.metric("Regulatory", herbal_info.get("Regulatory Risk", "—"))

            st.divider()
            st.markdown("### Market Positioning")
            st.write(f"**Market Type:** {herbal_info.get('Market Type', '—')}")

            # Optional expanded scientific layers if your compute_herbal_opportunity provides them
            if isinstance(herbal_info, dict) and "Predictive Toxicology" in herbal_info:
                st.divider()
                st.markdown("### Scientific Risk Layers")
                col1, col2 = st.columns(2)

                with col1:
                    tox = herbal_info.get("Predictive Toxicology", {})
                    st.markdown("#### Predictive Toxicology")
                    st.write(f"Risk Level: {tox.get('Risk Level', '—')}")
                    flags = tox.get("Flags", [])
                    if flags:
                        st.write("Flags:", ", ".join(flags))

                    st.markdown("#### Pharmacokinetics")
                    st.write(herbal_info.get("Pharmacokinetics", "—"))

                with col2:
                    st.markdown("#### Network Pharmacology")
                    st.write(herbal_info.get("Network Pharmacology", "—"))

                    st.markdown("#### Synergy Potential")
                    st.write(f"Score: {herbal_info.get('Synergy Potential Score', '—')}")

                    st.markdown("#### Plant Identification Risk")
                    st.write(herbal_info.get("Plant Identification Risk", "—"))

                    st.markdown("#### Quality Control Risk")
                    st.write(herbal_info.get("Quality Control Risk", "—"))

            overview_snapshot_lines += [
                f"Mode: Herbal",
                f"Herbal: {selected_drug}",
                f"Demand Score: {herbal_info.get('Demand Score', '—')}",
                f"Formulation Complexity: {herbal_info.get('Formulation Complexity', '—')}",
                f"Regulatory Risk: {herbal_info.get('Regulatory Risk', '—')}",
                f"Market Type: {herbal_info.get('Market Type', '—')}",
            ]
    else:
        if not selected_drug or not ob_data or ob_data.get("products") is None:
            st.info("Select a drug and click Analyze.")
            overview_snapshot_lines.append("No FDA drug selected or no Orange Book data.")
        else:
            # Executive overview
            try:
                overview = build_executive_overview(
                    drug_name=selected_drug,
                    route=st.session_state.get("route"),
                    ob_data=ob_data,
                    agent=st.session_state.get("meta_agent"),
                )
            except Exception as e:
                st.error(f"Failed to build executive overview: {e}")
                overview = {}

            formulation = infer_formulation_type(ob_data["products"]) or {
                "formulation_type": "Unknown",
                "risk_level": "Unknown",
            }

            formulation_risk = compute_formulation_risk(
                formulation=formulation,
                route=st.session_state.get("route"),
            )

            comp_count, comp_level = competitive_density_score(ob_data)

            ege_info = earliest_generic_entry(ob_data)
            ege_years = None
            if ege_info and ege_info.get("earliest_date") is not None:
                try:
                    ege_years = max(0.0, (ege_info["earliest_date"] - pd.Timestamp.today()).days / 365.0)
                except Exception:
                    ege_years = None

            cost = st.session_state.get("cost")

            c1, c2, c3, c4, c5 = st.columns(5)
            c1.metric("Drug", selected_drug)
            c2.metric("Formulation", formulation.get("formulation_type", "—"))
            c3.metric("Risk", formulation.get("risk_level", "—"))
            c4.metric("CMC Cost", f"${cost['low']}–${cost['high']}M" if cost else "—")
            c5.metric("Competition", f"{comp_level} ({comp_count})")

            st.divider()
            st.markdown("### Chemical Composition")
            api_name = "Unknown"
            chem = None
            try:
                if ob_data.get("products") is not None:
                    api_name = extract_api_from_orange_book(ob_data["products"])
                    chem = extract_chemical_profile(ob_data["products"])
            except Exception:
                chem = None

            if chem:
                st.write(f"**API:** {api_name}")
                st.write(f"**Salt/Form:** {chem.get('Salt / Form', 'Unknown')}")
                strength = (
                    chem.get("Strength")
                    or chem.get("strength")
                    or chem.get("strength_text")
                    or chem.get("strength_desc")
                    or "Unknown"
                )
                st.write(f"**Strength:** {strength}")
                st.write(f"**Dosage Form:** {chem.get('Dosage Form', 'Unknown')}")
                route_display = (
                    chem.get("Route")
                    or chem.get("route")
                    or chem.get("route_of_administration")
                    or st.session_state.get("route")
                    or "Unknown"
                )
                st.write(f"**Route:** {route_display}")
            else:
                st.write("Chemical composition unavailable.")

            # Recommendation
            flag, rationale = recommendation_flag(
                ege_years=ege_years,
                formulation_risk=formulation.get("risk_level"),
                density=comp_level,
                route=st.session_state.get("route"),
                cost=cost,
                competition=comp_level,
            )

            st.markdown("### Recommendation")
            {"Go": st.success, "Watch": st.warning, "Avoid": st.error}.get(flag, st.info)(
                f"{flag.upper()} — {rationale}"
            )

            if is_high_value_complex(formulation.get("risk_level"), comp_level, ege_years):
                st.info("💎 High-Value Complex Generic Opportunity Identified")
            cost_text = "Unknown" if not cost else f"${cost.get('low')}–${cost.get('high')}M"

            overview_snapshot_lines += [
                f"Mode: FDA Orange Book",
                f"Drug: {selected_drug}",
                f"Formulation: {formulation.get('formulation_type', '—')}",
                f"Risk: {formulation.get('risk_level', '—')}",
                f"Competition: {comp_level} ({comp_count})",
                f"EGE: {'Open now' if ege_years is None else f'{ege_years:.1f} years'}",
                f"CMC Cost: {cost_text}",
                f"Recommendation: {flag} — {rationale}",
            ]

    st.session_state.overview_snapshot = "\n".join(overview_snapshot_lines) if overview_snapshot_lines else ""


# =================================================
# EGE TAB
# =================================================
with ege_tab:
    st.subheader("⏳ Earliest Generic Entry (EGE) Dashboard")

    ege_snapshot_lines = []

    if commercial_mode != "FDA Orange Book (Rx Generics)":
        st.info("EGE applies only to FDA Orange Book drugs.")
        ege_snapshot_lines.append("Mode: Herbal — EGE not applicable.")
    else:
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("Total Products", len(products) if products is not None else 0)
        col2.metric("Generic-Ready Now", len(generic_df) if generic_df is not None else 0)
        col3.metric("Expiring ≤ 12 Months", len(expiring_df) if expiring_df is not None else 0)
        col4.metric("Tracked Opportunities", max((len(products) if products is not None else 0) - (len(generic_df) if generic_df is not None else 0), 0))

        ege_snapshot_lines += [
            f"Total Products: {len(products) if products is not None else 0}",
            f"Generic Ready: {len(generic_df) if generic_df is not None else 0}",
            f"Expiring ≤ 12 Months: {len(expiring_df) if expiring_df is not None else 0}",
        ]

        st.divider()
        st.markdown("### Selected Drug – EGE Detail")

        if not ob_data:
            st.info("Select a drug or patent to view Earliest Generic Entry.")
            ege_snapshot_lines.append("No selected drug EGE detail.")
        else:
            info = earliest_generic_entry(ob_data) or {}
            ege_date = info.get("earliest_date")

            if ege_date:
                st.metric("Earliest Generic Entry", ege_date.strftime("%B %d, %Y"))
                ege_snapshot_lines.append(f"Selected drug EGE date: {ege_date.strftime('%Y-%m-%d')}")
            else:
                st.success("No Orange Book blockers — generic eligible now")
                ege_snapshot_lines.append("Selected drug: Open now (no blockers).")

            with st.expander("Blocking Patents & Exclusivities"):
                patents_df = ob_data.get("patents")
                exclus_df = ob_data.get("exclusivities")

                if patents_df is not None and not patents_df.empty:
                    st.markdown("**Patents**")
                    st.dataframe(patents_df, use_container_width=True)
                else:
                    st.caption("No listed blocking patents.")

                if exclus_df is not None and not exclus_df.empty:
                    st.markdown("**Exclusivities**")
                    st.dataframe(exclus_df, use_container_width=True)
                else:
                    st.caption("No listed exclusivities.")

        st.divider()
        st.markdown("### Portfolio EGE Scanner")

        left, right = st.columns([1, 3])

        max_drugs = left.number_input(
            "Max drugs to scan",
            min_value=100,
            max_value=5000,
            value=800,
            step=100,
        )

        if left.button("Build / Refresh EGE Table", use_container_width=True):
            with st.spinner("Computing Earliest Generic Entry dates…"):
                try:
                    st.session_state.ege_master_df = build_ege_table_cached(int(max_drugs))
                except Exception as e:
                    st.error(f"Failed to build EGE table: {e}")

        ege_master_df = st.session_state.get("ege_master_df")

        if ege_master_df is None or getattr(ege_master_df, "empty", True):
            st.info("Click **Build / Refresh EGE Table** to generate portfolio EGE data.")
        else:
            f1, f2 = st.columns(2)
            max_years = f1.slider("Max Years to Entry", 0, 6, 3)
            status_filter = f2.multiselect("Status", ["Open", "Blocked"], ["Open", "Blocked"])

            filtered_df = (
                ege_master_df[
                    (ege_master_df["Years to Entry"] <= max_years)
                    & (ege_master_df["Status"].isin(status_filter))
                ]
                .sort_values("Years to Entry")
                .reset_index(drop=True)
            )

            st.dataframe(filtered_df, use_container_width=True, height=520)

    st.session_state.ege_snapshot = "\n".join(ege_snapshot_lines) if ege_snapshot_lines else ""


# =====================================================
# STRATEGY TAB
# =====================================================
with strategy_tab:
    st.subheader("📊 Market Strategy")

    strategy_snapshot_lines = []

    if commercial_mode != "FDA Orange Book (Rx Generics)":
        st.info("Strategy applies only to FDA Orange Book mode.")
        strategy_snapshot_lines.append("Mode: Herbal — strategy ranking not applicable.")
    else:
        st.markdown("### Generic-Ready Drugs (No Listed Blockers)")
        st.metric("Generic-ready products", len(generic_df) if generic_df is not None else 0)

        strategy_snapshot_lines.append(f"Generic-ready products: {len(generic_df) if generic_df is not None else 0}")

        if generic_df is None or generic_df.empty:
            st.info("No generic-ready drugs found.")
        else:
            st.dataframe(generic_df, use_container_width=True, height=420)

        st.divider()
        st.markdown("### Patents Expiring Soon")

        months_new = st.slider(
            "Months ahead",
            min_value=1,
            max_value=48,
            value=st.session_state.months_horizon,
            step=1,
        )

        if months_new != st.session_state.months_horizon:
            st.session_state.months_horizon = months_new

        expiring_df2 = pd.DataFrame(find_expiring_patents(patents, st.session_state.months_horizon))

        st.metric("Expiring patents", len(expiring_df2))

        strategy_snapshot_lines.append(f"Expiring patents (≤{st.session_state.months_horizon} months): {len(expiring_df2)}")

        if expiring_df2.empty:
            st.info("No patents expiring in the selected window.")
        else:
            st.dataframe(expiring_df2, use_container_width=True, height=420)

        st.divider()
        st.subheader("Top-20 Orange Book Opportunities (Auto-ranked)")

        run_rank = st.button("Run / Refresh Top-20 Ranking", use_container_width=True)

        if run_rank:
            with st.spinner("Ranking opportunities…"):
                try:
                    st.session_state.top20_opportunities = rank_top_orange_book_opportunities(
                        products=products,
                        patents=patents,
                        exclus=exclus,
                        top_n=20,
                        max_drugs=1500,
                    )
                    st.session_state.top20_computed = True
                except Exception as e:
                    st.error(f"Ranking failed: {e}")

        top20 = st.session_state.top20_opportunities

        if top20 is None:
            st.info("Click **Run / Refresh Top-20 Ranking** to generate opportunities.")
        elif getattr(top20, "empty", True):
            st.warning("No opportunities could be ranked from the current Orange Book data.")
        else:
            st.dataframe(top20, use_container_width=True, height=520)
            strategy_snapshot_lines.append("Top-20 ranking generated.")

    st.session_state.strategy_snapshot = "\n".join(strategy_snapshot_lines) if strategy_snapshot_lines else ""


# =====================================================
# AI TAB
# =====================================================
with ai_tab:
    st.subheader("🧠 AI Market Intelligence")

    if st.session_state.get("run_id") is None:
        st.info("Run Analyze first.")
    else:

        # =========================================
        # Build Context
        # =========================================
        if commercial_mode == "Herbal / Botanical Drugs":
            context_block = f"""
Herbal Product: {selected_drug}
Commercial Domain: Nutraceutical
"""
        else:
            context_block = f"""
Drug: {selected_drug}
Commercial Domain: FDA Orange Book Generic
"""

        # =========================================
        # Generate Initial AI Report
        # =========================================
        if st.button("Generate AI Executive Report"):

            base_prompt = f"""
You are a senior pharmaceutical commercialization strategist.

Context:
{context_block}

Provide:
1) Opportunity assessment
2) Key risks
3) Competitive positioning
4) Recommended next steps

No fabricated financial numbers.
Executive tone.
"""

            with st.spinner("Generating executive report..."):
                response = llm.invoke(base_prompt).content.strip()

            st.session_state.ai_summary = response
            if st.session_state.ai_chat_history:
               st.markdown("### 💬 AI Discussion")

            # Reverse order (latest first)
            for chat in reversed(st.session_state.ai_chat_history):
                st.markdown(f"**You:** {chat['user']}")
                st.markdown(f"**AI:** {chat['ai']}")
                st.markdown("---")


        # =========================================
        # Show Initial Report
        # =========================================
        if st.session_state.get("ai_summary"):
            st.markdown("### 📄 Executive AI Report")
            st.write(st.session_state.ai_summary)

        st.divider()

        # =========================================
        # AGENT — Ask Specific Question
        # =========================================
        st.markdown("### 🤖 Ask a Specific Question")

        user_question = st.text_input(
            "Ask about competition, toxicity, patent risk, launch timing, etc.",
            key="ai_followup_question",
        )

        if st.button("Ask AI") and user_question:

            # Build memory context
            memory_block = "\n".join(
                [f"User: {m['user']}\nAI: {m['ai']}" for m in st.session_state.ai_chat_history]
            )

            followup_prompt = f"""
You are a senior commercialization analyst.

Context:
{context_block}

Initial AI Report:
{st.session_state.get("ai_summary")}

Conversation History:
{memory_block}

User Question:
{user_question}

Answer specifically and concisely.
Do not invent financial numbers.
"""

            with st.spinner("Thinking..."):
                answer = llm.invoke(followup_prompt).content.strip()

            # Store memory
            st.session_state.ai_chat_history.append(
                {"user": user_question, "ai": answer}
            )

        # =========================================
        # Display Chat History
        # =========================================
        if st.session_state.ai_chat_history:
            st.markdown("### 💬 AI Discussion")

            for chat in st.session_state.ai_chat_history:
                st.markdown(f"**You:** {chat['user']}")
                st.markdown(f"**AI:** {chat['ai']}")
                st.markdown("---")

# =====================================================
# DATA UPLOAD TAB
# =====================================================
ORANGE_BOOK_DIR = "data/orange_book"

with data_upload:
    st.subheader("📂 Load Excel / CSV Orange Book Data")

    st.caption("Upload a dataset to overwrite local Orange Book CSVs used by the app.")

    uploaded_file = st.file_uploader(
        "Upload Excel or CSV",
        type=["xlsx", "xls", "csv"],
        accept_multiple_files=False,
        key="ob_uploader",
    )

    if uploaded_file:
        sheets = read_file_dynamic(uploaded_file)

        if not sheets:
            st.error("No readable sheets found.")
        else:
            sheet_names = list(sheets.keys())
            selected_sheet = st.selectbox("Select sheet", sheet_names, key="sheet_pick")

            sheet_info = sheets[selected_sheet]

            if "error" in sheet_info:
                st.error(sheet_info["error"])
            else:
                st.dataframe(sheet_info["data"], use_container_width=True)

                st.divider()
                st.markdown("### Promote as Orange Book Dataset")

                dataset_type = st.selectbox(
                    "Dataset type",
                    ["products", "patents", "exclusivities"],
                    key="dataset_type",
                )

                if st.button("Overwrite & Reload App", type="primary", key="overwrite_reload"):
                    os.makedirs(ORANGE_BOOK_DIR, exist_ok=True)

                    target_path = os.path.join(ORANGE_BOOK_DIR, f"{dataset_type}.csv")
                    full_df = sheet_info["data"]
                    full_df.to_csv(target_path, index=False)

                    st.success(f"✅ `{dataset_type}.csv` updated at {target_path}")
                    reload_app_after_upload()


# =====================================================
# Export Section (Bottom)
# =====================================================
st.divider()
st.subheader("📤 Export Report")

# Build PPT content from stored snapshots
ppt_overview = st.session_state.get("overview_snapshot", "") or "—"
ppt_ege = st.session_state.get("ege_snapshot", "") or "—"
ppt_strategy = st.session_state.get("strategy_snapshot", "") or "—"
ppt_ai = st.session_state.get("ai_summary") or st.session_state.get("ai_snapshot", "") or "—"

colA, colB = st.columns([1, 2])
with colA:
    export_btn = st.button("⬇ Build PowerPoint", use_container_width=True)

with colB:
    st.caption("Exports Overview, EGE, Strategy, and full AI text (with bullet auto-detection).")

if export_btn:
    ppt_file = export_report_to_ppt(
        commercial_mode=commercial_mode,
        selected_drug=selected_drug,
        overview_text=ppt_overview,
        ege_text=ppt_ege,
        strategy_text=ppt_strategy,
        ai_text=ppt_ai,
    )

    st.download_button(
        label="Download PowerPoint",
        data=ppt_file,
        file_name=f"{(selected_drug or 'Report').replace(' ', '_')}_Commercial_Report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )
